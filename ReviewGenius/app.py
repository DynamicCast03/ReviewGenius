import os
import json
import pdfplumber # 导入 pdfplumber
import pptx # 导入 pptx
from flask import (
    Flask,
    request,
    Response,
    render_template,
    jsonify,
)
from flask_socketio import SocketIO, emit
from openai import AuthenticationError
# from werkzeug.utils import secure_filename
from filter import sanitizer
import prompt_manager
import siliconflow_client
from llm_json_parser import stream_json_with_events
from question_types import Question # Import base class for validation
import grading # 导入评分模块

app = Flask(__name__)
# 添加CORS和SocketIO
socketio = SocketIO(app, cors_allowed_origins="*")

UPLOAD_FOLDER = "uploads"
CONFIG_FILE = "config.json"
app.config["MAX_CONTENT_LENGTH"] = 50 * 1024 * 1024
app.config["UPLOAD_EXTENSIONS"] = [".txt", ".pdf", ".pptx", ".ppt"]  # 支持txt, pdf, pptx, ppt
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

def get_uploaded_files():
    """获取上传目录中的文件列表，忽略隐藏文件。"""
    try:
        # Ignore hidden files like .DS_Store
        return [f for f in os.listdir(UPLOAD_FOLDER) if os.path.isfile(os.path.join(UPLOAD_FOLDER, f)) and not f.startswith('.')]
    except Exception:
        return []

def broadcast_file_list():
    """向所有连接的客户端广播当前的文件列表"""
    with app.app_context():
        files = get_uploaded_files()
        socketio.emit('file_list_update', {'files': files})

def load_config():
    if not os.path.exists(CONFIG_FILE):
        default_config = {
            "temperature": 1.0,
            "enhanced_structured_output": False
        }
        with open(CONFIG_FILE, 'w', encoding='utf-8') as f:
            json.dump(default_config, f, indent=4)
        return default_config
    try:
        with open(CONFIG_FILE, 'r', encoding='utf-8') as f:
            return json.load(f)
    except (json.JSONDecodeError, FileNotFoundError):
        return {
            "temperature": 1.0,
            "enhanced_structured_output": False
        }

def save_config(config_data):
    with open(CONFIG_FILE, 'w', encoding='utf-8') as f:
        json.dump(config_data, f, indent=4)

@app.route("/")
def index():
    config = load_config()
    return render_template("./interface.html", config=config)

@app.route("/api/settings", methods=["GET", "POST"])
def manage_settings():
    if request.method == "GET":
        return jsonify(load_config())
    
    if request.method == "POST":
        data = request.json
        config = load_config()
        config["temperature"] = float(data.get("temperature", config["temperature"]))
        config["enhanced_structured_output"] = bool(data.get("enhanced_structured_output", config["enhanced_structured_output"]))
        save_config(config)
        return jsonify({"message": "设置已保存", "config": config})

@app.route("/api/upload", methods=["POST"])
def upload_file():
    if 'files' not in request.files:
        return jsonify({"error": "没有文件部分"}), 400
    
    uploaded_files = request.files.getlist("files")
    uploaded_files = [f for f in uploaded_files if f.filename]

    if not uploaded_files:
        return jsonify({"error": "未上传或未选择任何文件"}), 400

    errors = {}
    success_files = []
    for file in uploaded_files:
        # filename = secure_filename(file.filename)
        filename = file.filename
        file_ext = os.path.splitext(filename)[1].lower()
        if file_ext not in app.config["UPLOAD_EXTENSIONS"]:
            errors[file.filename] = f"不支持的文件类型"
            continue
        
        try:
            file.save(os.path.join(UPLOAD_FOLDER, filename))
            success_files.append(filename)
        except Exception as e:
            errors[file.filename] = f"保存文件失败: {str(e)}"

    if not errors:
        broadcast_file_list()
        return jsonify({"message": "文件上传成功"}), 200

    # 如果有错误，也广播列表以反映部分成功的上传
    broadcast_file_list()
    response = {"message": "部分或全部文件上传失败", "errors": errors, "success_files": success_files}
    return jsonify(response), 400 if len(success_files) == 0 else 207

@app.route("/api/files", methods=["GET"])
def list_files():
    try:
        files = get_uploaded_files()
        response = jsonify(files)
        # 防止浏览器缓存文件列表
        response.headers['Cache-Control'] = 'no-cache, no-store, must-revalidate'
        response.headers['Pragma'] = 'no-cache'
        response.headers['Expires'] = '0'
        return response
    except Exception as e:
        return jsonify({"error": f"无法列出文件: {str(e)}"}), 500

@app.route("/api/files/<filename>", methods=["DELETE"])
def delete_file_route(filename):
    # werkzeug.utils.secure_filename() is used to secure the filename before saving it.
    # We should use the same function to secure the filename before deleting it.
    # filename = secure_filename(filename)
    file_path = os.path.join(UPLOAD_FOLDER, filename)
    
    if not os.path.exists(file_path):
        return jsonify({"error": "文件未找到"}), 404
        
    try:
        os.remove(file_path)
        broadcast_file_list() # 广播更新后的文件列表
        return jsonify({"message": f"文件 '{filename}' 已删除"}), 200
    except Exception as e:
        return jsonify({"error": f"删除文件失败: {str(e)}"}), 500

@app.route("/api/process", methods=["POST"])
def generate_exam():
    # 更改: 不再从请求中获取文件，而是扫描上传文件夹
    uploaded_filenames = [f for f in os.listdir(UPLOAD_FOLDER) if os.path.isfile(os.path.join(UPLOAD_FOLDER, f)) and not f.startswith('.')]

    if not uploaded_filenames:
        return jsonify({"error": "请先上传至少一个参考资料文件"}), 400

    try:
        user_text = request.form.get("user_input", "无特定要求")
        user_text = sanitizer.sanitize(user_text)
        api_key = request.form.get("api_key")
        
        config = load_config()
        temperature = config.get("temperature", 1.0)

        if not api_key:
            return jsonify({"error": "API Key缺失"}), 400

        # 合并简答题和计算题
        short_answer_count = int(request.form.get("short_count", "0")) + int(
            request.form.get("calc_count", "0")
        )

        question_settings = {
            "选择题": {
                "count": request.form.get("choice_count", "0"),
                "score": request.form.get("choice_score", "5"),
            },
            "填空题": {
                "count": request.form.get("blank_count", "0"),
                "score": request.form.get("blank_score", "5"),
            },
            "简答题": {
                "count": str(short_answer_count),
                "score": request.form.get("short_score", "10"),
            },
        }

        question_types_str = "、".join(
            [
                f"{k}{v['count']}道(每题{v['score']}分)"
                for k, v in question_settings.items()
                if int(v['count']) > 0
            ]
        )
        if not question_types_str:
            return jsonify({"error": "至少需要设置一种题型"}), 400

        document_contents = []
        # 更改: 循环读取文件夹中的文件
        for filename in uploaded_filenames:
            file_path = os.path.join(UPLOAD_FOLDER, filename)
            file_ext = os.path.splitext(filename)[1].lower()
            content = ""
            try:
                if file_ext == ".pdf":
                    with pdfplumber.open(file_path) as pdf:
                        all_text = [page.extract_text() for page in pdf.pages if page.extract_text()]
                        content = "\\n".join(all_text)
                elif file_ext in [".pptx", ".ppt"]:
                    pres = pptx.Presentation(file_path)
                    all_text = []
                    for slide in pres.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                all_text.append(shape.text)
                    content = "\\n".join(all_text)
                else: # 默认为txt
                    with open(file_path, 'r', encoding='utf-8') as f:
                        content = f.read()
            except Exception as e:
                 return jsonify({"error": f"读取文件 '{filename}' 时出错: {str(e)}"}), 500

            document_contents.append(f"--- 来自文件: {filename} ---\n{content}")
        document_content = "\\n\\n".join(document_contents)

        scores_data = {
            "multiple_choice": question_settings["选择题"]["score"],
            "fill_in_the_blank": question_settings["填空题"]["score"],
            "short_answer": question_settings["简答题"]["score"],
        }

        # 1. 加载静态的格式化提示词
        formatting_instructions = prompt_manager.get_prompt("exam_generation_prompt_formatting")

        # 2. 渲染主提示词，注入所有动态内容和静态格式说明
        main_prompt = prompt_manager.get_prompt("exam_generation_prompt", document_content=document_content,
        user_requirement=user_text,
        question_types=question_types_str,
        formatting_instructions=formatting_instructions,
        scores=scores_data)
        
        messages = [{"role": "user", "content": main_prompt}]

        def generate_question_stream():
            try:
                # 检查是否启用增强模式
                config = load_config()
                enhanced_mode = config.get("enhanced_structured_output", False)

                # 使用我们创建的客户端进行流式调用
                llm_stream = siliconflow_client.invoke_llm(
                    api_key=api_key,
                    model="Qwen/Qwen2.5-72B-Instruct",
                    messages=messages,
                    stream=not enhanced_mode, # 如果增强模式，第一次调用就不是流式
                    temperature=temperature,
                    enhanced_structured_output=enhanced_mode,
                    formatting_prompt=formatting_instructions if enhanced_mode else None
                )
                
                # 如果是增强模式，llm_stream已经是第二次调用的流，直接处理
                if enhanced_mode:
                    event_stream = stream_json_with_events(llm_stream)
                    for event in event_stream:
                        yield json.dumps(event) + "\\n"
                    return

                # 原有逻辑，处理非增强模式下的流
                event_stream = stream_json_with_events(llm_stream)

                for event in event_stream:
                    if event["type"] == "end":
                        # 在结束后验证数据结构
                        try:
                            question_obj = Question.from_dict(event["data"])
                            # 将验证和转换后的数据放回事件中
                            event["data"] = question_obj.to_dict()
                        except (ValueError, KeyError) as e:
                            # 如果数据格式错误，可以跳过或发送一个错误事件
                            print(f"Skipping invalid question object: {e}, data: {event['data']}")
                            continue # 不发送这个 'end' 事件
                    
                    yield json.dumps(event) + "\\n"

            except AuthenticationError:
                yield json.dumps({"type": "error", "error": "API Key 无效或已过期，请检查您的输入。", "error_type": "authentication"}) + "\\n"
            except ValueError as e:
                # 捕获由内容安全检查抛出的ValueError
                if "输入内容被判定为不安全" in str(e):
                    yield json.dumps({"type": "error", "error": str(e), "error_type": "security"}) + "\\n"
                else:
                    # 重新抛出其他类型的ValueError
                    yield json.dumps({"type": "error", "error": f"生成过程中发生验证错误: {str(e)}", "error_type": "generation"}) + "\\n"
            except Exception as e:
                # 捕获其他流式过程中的错误
                yield json.dumps({"type": "error", "error": f"生成过程中发生错误: {str(e)}", "error_type": "generation"}) + "\\n"

        return Response(generate_question_stream(), mimetype="application/x-ndjson")

    except Exception as e:
        error_msg = f"处理时出错: {str(e)}"
        print(error_msg)
        return jsonify({"error": error_msg}), 500


@app.route("/api/grade", methods=["POST"])
def grade_submission():
    try:
        data = request.get_json()
        questions = data.get("questions")
        user_answers = data.get("answers")
        api_key = data.get("api_key")
        
        config = load_config()
        temperature = config.get("temperature", 0.7)
        enhanced_mode = config.get("enhanced_structured_output", False)

        if not all([questions, user_answers, api_key]):
            return jsonify({"error": "缺少题目、答案或API Key"}), 400

        def generate_grade_stream():
            try:
                grading_stream = grading.grade_exam_stream(
                    questions, user_answers, api_key, temperature,
                    enhanced_structured_output=enhanced_mode
                )
                for event in grading_stream:
                    yield event
            except Exception as e:
                error_event = {
                    "type": "error",
                    "error": f"启动批改流失败: {str(e)}",
                }
                yield json.dumps(error_event) + "\\n"

        return Response(generate_grade_stream(), mimetype="application/x-ndjson")

    except Exception as e:
        error_msg = f"评分接口出错: {str(e)}"
        print(error_msg)
        return jsonify({"error": error_msg}), 500

@socketio.on('connect')
def handle_connect():
    """当客户端连接时，立即向其发送当前的文件列表"""
    print('Client connected')
    # 使用 with app.app_context() 来确保可以访问应用上下文
    with app.app_context():
        files = get_uploaded_files()
        emit('file_list_update', {'files': files})

@socketio.on('disconnect')
def handle_disconnect():
    print('Client disconnected')

if __name__ == "__main__":
    socketio.run(app, host="0.0.0.0", port=5000, debug=True)
