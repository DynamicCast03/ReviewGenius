import os
import json
import logging
from logging.handlers import RotatingFileHandler
import pdfplumber 
import pptx 
from flask import (
    Flask,
    request,
    Response,
    render_template,
    jsonify,
)
from flask_socketio import SocketIO, emit
from openai import AuthenticationError

from filter import sanitizer
import prompt_manager
import siliconflow_client
from llm_json_parser import stream_json_with_events
from question_types import Question 
import grading 
import markdown_exporter 
import user_profile_manager
import threading

def setup_logging():
    log_dir = 'log'
    if not os.path.exists(log_dir):
        os.makedirs(log_dir)
    
    log_file = os.path.join(log_dir, 'app.log')
    
    handler = RotatingFileHandler(log_file, maxBytes=10*1024*1024, backupCount=5, encoding='utf-8')
    stream_handler = logging.StreamHandler()
    
    formatter = logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(message)s')
    handler.setFormatter(formatter)
    stream_handler.setFormatter(formatter)
    
    logger = logging.getLogger()
    logger.setLevel(logging.INFO)
    logger.addHandler(handler)
    logger.addHandler(stream_handler)

setup_logging()

app = Flask(__name__)
socketio = SocketIO(app, cors_allowed_origins="*")

UPLOAD_FOLDER = "uploads"
CONFIG_FILE = "config.json"
app.config["MAX_CONTENT_LENGTH"] = 50 * 1024 * 1024
app.config["UPLOAD_EXTENSIONS"] = [".txt", ".pdf", ".pptx", ".ppt"]  
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

def get_uploaded_files():
    """获取上传目录中的文件列表，忽略隐藏文件。"""
    try:
        return [f for f in os.listdir(UPLOAD_FOLDER) if os.path.isfile(os.path.join(UPLOAD_FOLDER, f)) and not f.startswith('.')]
    except Exception:
        return []

def broadcast_file_list():
    """向所有连接的客户端广播当前的文件列表"""
    with app.app_context():
        files = get_uploaded_files()
        socketio.emit('file_list_update', {'files': files})

def load_config():
    default_profile = "该用户暂无画像，请根据本次答题情况生成一份初始画像。"
    if not os.path.exists(CONFIG_FILE):
        app.logger.info(f"配置文件不存在，创建默认配置文件: {CONFIG_FILE}")
        default_config = {
            "temperature": 1.0,
            "enhanced_structured_output": False,
            "user_profile": default_profile,
            "user_profile_enabled": True
        }
        with open(CONFIG_FILE, 'w', encoding='utf-8') as f:
            json.dump(default_config, f, indent=4, ensure_ascii=False)
        return default_config
    try:
        app.logger.info(f"加载配置文件: {CONFIG_FILE}")
        with open(CONFIG_FILE, 'r', encoding='utf-8') as f:
            config = json.load(f)
            if "user_profile" not in config:
                config["user_profile"] = default_profile
            if "user_profile_enabled" not in config:
                config["user_profile_enabled"] = True
            return config
    except (json.JSONDecodeError, FileNotFoundError):
        app.logger.error(f"加载配置文件失败: {CONFIG_FILE}")
        return {
            "temperature": 1.0,
            "enhanced_structured_output": False,
            "user_profile": default_profile,
            "user_profile_enabled": True
        }

def save_config(config_data):
    app.logger.info(f"保存配置文件: {CONFIG_FILE}")
    with open(CONFIG_FILE, 'w', encoding='utf-8') as f:
        json.dump(config_data, f, indent=4, ensure_ascii=False)

@app.route("/")
def index():
    config = load_config()
    return render_template("./interface.html", config=config)

@app.route("/api/settings", methods=["GET", "POST"])
def manage_settings():
    if request.method == "GET":
        app.logger.info(f"获取配置文件: {CONFIG_FILE}")
        return jsonify(load_config())
    
    if request.method == "POST":
        app.logger.info(f"更新配置文件: {CONFIG_FILE}")
        data = request.json
        config = load_config()
        config["temperature"] = float(data.get("temperature", config["temperature"]))
        config["enhanced_structured_output"] = bool(data.get("enhanced_structured_output", config["enhanced_structured_output"]))
        config["user_profile_enabled"] = bool(data.get("user_profile_enabled", config.get("user_profile_enabled", True)))
        if "user_profile" in data:
            config["user_profile"] = str(data.get("user_profile", ""))
        
        save_config(config)
        return jsonify({"message": "设置已保存", "config": config})

@app.route("/api/upload", methods=["POST"])
def upload_file():
    app.logger.info(f"上传文件: {UPLOAD_FOLDER}")
    if 'files' not in request.files:
        return jsonify({"error": "没有文件部分"}), 400
    
    uploaded_files = request.files.getlist("files")
    uploaded_files = [f for f in uploaded_files if f.filename]

    if not uploaded_files:
        return jsonify({"error": "未上传或未选择任何文件"}), 400

    errors = {}
    success_files = []
    for file in uploaded_files:
        filename = file.filename
        file_ext = os.path.splitext(filename)[1].lower()
        if file_ext not in app.config["UPLOAD_EXTENSIONS"]:
            errors[file.filename] = f"不支持的文件类型"
            continue
        
        try:
            file.save(os.path.join(UPLOAD_FOLDER, filename))
            success_files.append(filename)
        except Exception as e:
            errors[file.filename] = f"保存文件失败: {str(e)}"

    if not errors:
        broadcast_file_list()
        return jsonify({"message": "文件上传成功"}), 200

    broadcast_file_list()
    response = {"message": "部分或全部文件上传失败", "errors": errors, "success_files": success_files}
    return jsonify(response), 400 if len(success_files) == 0 else 207

@app.route("/api/files", methods=["GET"])
def list_files():
    try:
        files = get_uploaded_files()
        response = jsonify(files)
        response.headers['Cache-Control'] = 'no-cache, no-store, must-revalidate'
        response.headers['Pragma'] = 'no-cache'
        response.headers['Expires'] = '0'
        return response
    except Exception as e:
        return jsonify({"error": f"无法列出文件: {str(e)}"}), 500

@app.route("/api/files/<filename>", methods=["DELETE"])
def delete_file_route(filename):
    file_path = os.path.join(UPLOAD_FOLDER, filename)
    
    if not os.path.exists(file_path):
        return jsonify({"error": "文件未找到"}), 404
        
    try:
        os.remove(file_path)
        broadcast_file_list()
        return jsonify({"message": f"文件 '{filename}' 已删除"}), 200
    except Exception as e:
        return jsonify({"error": f"删除文件失败: {str(e)}"}), 500

@app.route("/api/process", methods=["POST"])
def generate_exam():
    app.logger.info(f"开始生成试卷.")
    uploaded_filenames = [f for f in os.listdir(UPLOAD_FOLDER) if os.path.isfile(os.path.join(UPLOAD_FOLDER, f)) and not f.startswith('.')]

    if not uploaded_filenames:
        return jsonify({"error": "请先上传至少一个参考资料文件"}), 400

    try:
        user_text = request.form.get("user_input", "无特定要求")
        user_text = sanitizer.sanitize(user_text)
        api_key = request.form.get("api_key")
        
        config = load_config()
        temperature = config.get("temperature", 1.0)
        
        if config.get("user_profile_enabled", True):
            user_profile = config.get("user_profile", "该用户暂无画像，请根据本次答题情况生成一份初始画像。")
        else:
            user_profile = "用户画像功能未开启。"

        if not api_key:
            return jsonify({"error": "API Key缺失"}), 400

        short_answer_count = int(request.form.get("short_count", "0")) + int(
            request.form.get("calc_count", "0")
        )

        question_settings = {
            "选择题": {
                "count": request.form.get("choice_count", "0"),
                "score": request.form.get("choice_score", "5"),
            },
            "填空题": {
                "count": request.form.get("blank_count", "0"),
                "score": request.form.get("blank_score", "5"),
            },
            "简答题": {
                "count": str(short_answer_count),
                "score": request.form.get("short_score", "10"),
            },
        }

        question_types_str = "、".join(
            [
                f"{k}{v['count']}道(每题{v['score']}分)"
                for k, v in question_settings.items()
                if int(v['count']) > 0
            ]
        )
        if not question_types_str:
            return jsonify({"error": "至少需要设置一种题型"}), 400

        app.logger.info(f"开始生成试卷. 文件: {uploaded_filenames}, 要求: {question_types_str}")

        document_contents = []
        for filename in uploaded_filenames:
            file_path = os.path.join(UPLOAD_FOLDER, filename)
            file_ext = os.path.splitext(filename)[1].lower()
            content = ""
            try:
                if file_ext == ".pdf":
                    with pdfplumber.open(file_path) as pdf:
                        all_text = [page.extract_text() for page in pdf.pages if page.extract_text()]
                        content = "\n".join(all_text)
                elif file_ext in [".pptx", ".ppt"]:
                    pres = pptx.Presentation(file_path)
                    all_text = []
                    for slide in pres.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                all_text.append(shape.text)
                    content = "\n".join(all_text)
                else:
                    with open(file_path, 'r', encoding='utf-8') as f:
                        content = f.read()
            except Exception as e:
                 return jsonify({"error": f"读取文件 '{filename}' 时出错: {str(e)}"}), 500

            document_contents.append(f"--- 来自文件: {filename} ---\n{content}")
        document_content = "\n\n".join(document_contents)
        app.logger.info(f"所有文件内容已聚合，总长度: {len(document_content)}")

        scores_data = {
            "multiple_choice": question_settings["选择题"]["score"],
            "fill_in_the_blank": question_settings["填空题"]["score"],
            "short_answer": question_settings["简答题"]["score"],
        }

        formatting_instructions = prompt_manager.get_prompt("exam_generation_prompt_formatting")

        main_prompt = prompt_manager.get_prompt("exam_generation_prompt", document_content=document_content,
        user_requirement=user_text,
        question_types=question_types_str,
        formatting_instructions=formatting_instructions,
        scores=scores_data,
        user_profile=user_profile)
        
        messages = [{"role": "user", "content": main_prompt}]

        def generate_question_stream():
            try:
                config = load_config()
                enhanced_mode = config.get("enhanced_structured_output", False)

                llm_stream = siliconflow_client.invoke_llm(
                    api_key=api_key,
                    model="Qwen/Qwen2.5-72B-Instruct",
                    messages=messages,
                    stream=True,
                    temperature=temperature,
                    enhanced_structured_output=enhanced_mode,
                    formatting_prompt=formatting_instructions if enhanced_mode else None
                )
                
                if enhanced_mode:
                    event_stream = stream_json_with_events(llm_stream)
                    for event in event_stream:
                        yield json.dumps(event) + "\n"
                    return

                event_stream = stream_json_with_events(llm_stream)

                for event in event_stream:
                    if event["type"] == "end":
                        try:
                            question_obj = Question.from_dict(event["data"])
                            event["data"] = question_obj.to_dict()
                        except (ValueError, KeyError) as e:
                            app.logger.warning(f"Skipping invalid question object: {e}, data: {event['data']}")
                            continue
                    
                    yield json.dumps(event) + "\n"

            except AuthenticationError:
                yield json.dumps({"type": "error", "error": "API Key 无效或已过期，请检查您的输入。", "error_type": "authentication"}) + "\n"
            except ValueError as e:
                if "输入内容被判定为不安全" in str(e):
                    yield json.dumps({"type": "error", "error": str(e), "error_type": "security"}) + "\n"
                else:
                    yield json.dumps({"type": "error", "error": f"生成过程中发生验证错误: {str(e)}", "error_type": "generation"}) + "\n"
            except Exception as e:
                yield json.dumps({"type": "error", "error": f"生成过程中发生错误: {str(e)}", "error_type": "generation"}) + "\n"

        app.logger.info("返回试卷生成流.")
        return Response(generate_question_stream(), mimetype="application/x-ndjson")

    except Exception as e:
        error_msg = f"处理时出错: {str(e)}"
        app.logger.error(error_msg)
        return jsonify({"error": error_msg}), 500

@app.route("/api/regenerate_question", methods=["POST"])
def regenerate_question():
    try:
        data = request.get_json()
        original_question = data.get("question")
        action = data.get("action") 
        user_requirement = data.get("user_requirement", "无特定要求")
        api_key = data.get("api_key")
        
        app.logger.info(f"开始题目再生成. Action: {action}, Q_Type: {original_question.get('question_type')}")

        if not all([original_question, action, api_key]):
            return jsonify({"error": "缺少原始题目、操作类型或API Key"}), 400

        uploaded_filenames = get_uploaded_files()
        if not uploaded_filenames:
            return jsonify({"error": "找不到参考资料文件"}), 400

        document_contents = []
        for filename in uploaded_filenames:
            file_path = os.path.join(UPLOAD_FOLDER, filename)
            file_ext = os.path.splitext(filename)[1].lower()
            content = ""
            try:
                if file_ext == ".pdf":
                    with pdfplumber.open(file_path) as pdf:
                        all_text = [page.extract_text() for page in pdf.pages if page.extract_text()]
                        content = "\n".join(all_text)
                elif file_ext in [".pptx", ".ppt"]:
                    pres = pptx.Presentation(file_path)
                    all_text = [shape.text for slide in pres.slides for shape in slide.shapes if hasattr(shape, "text")]
                    content = "\n".join(all_text)
                else:
                    with open(file_path, 'r', encoding='utf-8') as f:
                        content = f.read()
            except Exception as e:
                 return jsonify({"error": f"读取文件 '{filename}' 时出错: {str(e)}"}), 500
            document_contents.append(f"--- 来自文件: {filename} ---\n{content}")
        
        document_content = "\n\n".join(document_contents)
        app.logger.info(f"题目再生成 - 文件内容已聚合，总长度: {len(document_content)}")

        prompt_map = {
            "regenerate": "regenerate_question_prompt",
            "increase_difficulty": "increase_difficulty_prompt",
            "decrease_difficulty": "decrease_difficulty_prompt",
        }
        prompt_name = prompt_map.get(action)
        if not prompt_name:
            return jsonify({"error": "无效的操作类型"}), 400

        config = load_config()
        temperature = config.get("temperature", 1.0)
        formatting_instructions = prompt_manager.get_prompt("exam_generation_prompt_formatting")

        main_prompt = prompt_manager.get_prompt(
            prompt_name,
            document_content=document_content,
            user_requirement=user_requirement,
            original_question=json.dumps(original_question, ensure_ascii=False, indent=2),
            score=original_question.get('score', 5), 
            formatting_instructions=formatting_instructions
        )

        messages = [{"role": "user", "content": main_prompt}]

        def generate_stream():
            try:
                enhanced_mode = config.get("enhanced_structured_output", False)
                llm_stream = siliconflow_client.invoke_llm(
                    api_key=api_key,
                    model="Qwen/Qwen2.5-72B-Instruct",
                    messages=messages,
                    stream=True,
                    temperature=temperature,
                    enhanced_structured_output=enhanced_mode,
                    formatting_prompt=formatting_instructions if enhanced_mode else None
                )

                event_stream = stream_json_with_events(llm_stream)
                for event in event_stream:
                    if event["type"] == "end":
                        try:
                            question_obj = Question.from_dict(event["data"])
                            question_obj.score = original_question.get('score', 5)
                            event["data"] = question_obj.to_dict()
                        except (ValueError, KeyError) as e:
                            app.logger.warning(f"Skipping invalid regenerated question object: {e}, data: {event['data']}")
                            continue
                    yield json.dumps(event) + "\n"

            except AuthenticationError:
                yield json.dumps({"type": "error", "error": "API Key 无效或已过期。", "error_type": "authentication"}) + "\n"
            except ValueError as e:
                 yield json.dumps({"type": "error", "error": str(e), "error_type": "security"}) + "\n"
            except Exception as e:
                yield json.dumps({"type": "error", "error": f"生成过程中发生错误: {str(e)}", "error_type": "generation"}) + "\n"

        app.logger.info("返回题目再生成流.")
        return Response(generate_stream(), mimetype="application/x-ndjson")

    except Exception as e:
        error_msg = f"处理时出错: {str(e)}"
        app.logger.error(error_msg)
        return jsonify({"error": error_msg}), 500

@app.route("/api/export/markdown", methods=["POST"])
def export_markdown():
    try:
        data = request.get_json()
        questions = data.get("questions")
        answers_placement = data.get("answers_placement", "inline")

        if not questions:
            return jsonify({"error": "没有提供题目数据"}), 400

        markdown_content = markdown_exporter.export_to_markdown(questions, answers_placement)
        
        return jsonify({"markdown_content": markdown_content})

    except Exception as e:
        error_msg = f"导出Markdown时出错: {str(e)}"
        app.logger.error(error_msg)
        return jsonify({"error": error_msg}), 500

def _update_profile_task(api_key, questions, user_answers):
    """在后台线程中生成答题总结并更新用户画像。"""
    with app.app_context():
        app.logger.info("后台任务：开始更新用户画像。")
        try:
            
            summary_prompt = prompt_manager.get_prompt(
                "grading_summary_prompt",
                questions_with_answers=json.dumps(questions, ensure_ascii=False, indent=2),
                user_answers=json.dumps(user_answers, ensure_ascii=False, indent=2),
            )
            
            
            llm_response = siliconflow_client.invoke_llm(
                api_key=api_key,
                model="Qwen/Qwen2.5-72B-Instruct",
                messages=[{"role": "user", "content": summary_prompt}],
                stream=False,
                temperature=0.6 
            )
            grading_summary = llm_response.choices[0].message.content.strip()
            app.logger.info(f"后台任务：生成答题总结完成，长度: {len(grading_summary)}")

            if not grading_summary:
                app.logger.info("LLM返回了空的答题总结，跳过用户画像更新。")
                return

            
            app.logger.info(f"生成的答题总结: {grading_summary}")
            updated_profile = user_profile_manager.update_user_profile(grading_summary, api_key)

            if updated_profile:
                app.logger.info(f"用户画像已成功更新: {updated_profile}")
            else:
                app.logger.warning("用户画像更新失败。")

        except Exception as e:
            app.logger.error(f"后台更新用户画像任务失败: {e}")

@app.route("/api/grade", methods=["POST"])
def grade_submission():
    try:
        data = request.get_json()
        questions = data.get("questions")
        user_answers = data.get("answers")
        api_key = data.get("api_key")
        
        config = load_config()
        temperature = config.get("temperature", 0.7)
        enhanced_mode = config.get("enhanced_structured_output", False)

        if not all([questions, user_answers, api_key]):
            return jsonify({"error": "缺少题目、答案或API Key"}), 400

        app.logger.info(f"开始批改试卷. 题目数: {len(questions)}")

        def generate_grade_stream():
            grading_results = []
            try:
                grading_stream = grading.grade_exam_stream(
                    questions, user_answers, api_key, temperature,
                    enhanced_structured_output=enhanced_mode
                )
                for event_str in grading_stream:
                    yield event_str + "\n" 
                
                
                config = load_config()
                if config.get("user_profile_enabled", True):
                    app.logger.info("用户画像功能已启用，启动后台任务更新用户画像。")
                    thread = threading.Thread(
                        target=_update_profile_task,
                        args=(api_key, questions, user_answers)
                    )
                    thread.start()
                else:
                    app.logger.info("用户画像功能未启用，跳过更新。")

            except Exception as e:
                error_event = {
                    "type": "error",
                    "error": f"启动批改流失败: {str(e)}",
                }
                yield json.dumps(error_event) + "\n"

        app.logger.info("返回批改结果流.")
        return Response(generate_grade_stream(), mimetype="application/x-ndjson")

    except Exception as e:
        error_msg = f"评分接口出错: {str(e)}"
        app.logger.error(error_msg)
        return jsonify({"error": error_msg}), 500

@socketio.on('connect')
def handle_connect():
    """当客户端连接时，立即向其发送当前的文件列表"""
    app.logger.info('Client connected')
    
    with app.app_context():
        files = get_uploaded_files()
        emit('file_list_update', {'files': files})

@socketio.on('disconnect')
def handle_disconnect():
    app.logger.info('Client disconnected')

if __name__ == "__main__":
    socketio.run(app, host="0.0.0.0", port=5000, debug=True)
